#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Extract workshop content into course-ta-agent/mcp-server-py/data/.

來源：
  - 6 個 .pptx（用 python-pptx 抽 text + table cells）
  - 已是 md 的（04-hands-on-lab.md / 05-practical-considerations.md / Lab 手冊 / benchmark）→ 直接 cp

輸出：course-ta-agent/mcp-server-py/data/*.md
  - 01-why-mcp.md
  - 02-how-mcp-works.md
  - 03-agentic-tool-loop.md
  - 04-hands-on-lab.md          （直接 cp 自 root）
  - 04-hands-on-lab-slides.md   （從 04-hands-on-lab.pptx 抽）
  - 05-practical-considerations.md（直接 cp）
  - lab-L1.md / lab-L2.md / lab-L3.md（cp 自 mini-project/docs/labs/）
  - benchmark-claude-vs-gemma4.md（cp 自 mini-project/docs/benchmarks/）
  - haiku-alignment-report.md   （直接 cp）
  - haiku-alignment-slides.md   （從 haiku-alignment-report.pptx 抽）
  - sonnet-running-example-slides.md（從 sonnet-running-example.pptx 抽）

用法：
  uv run --with python-pptx python3 tools/extract-pptx-to-md.py

每次 .pptx 修改後重跑一次。
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parent.parent
DATA_DIR = REPO_ROOT / "course-ta-agent" / "mcp-server-py" / "data"


def extract_pptx(pptx_path: Path) -> str:
    """抽 pptx 全部文字（含 table cells），每張 slide 用 --- 分隔。"""
    prs = Presentation(pptx_path)
    out: list[str] = [f"# {pptx_path.name}\n"]
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n## Slide {i}\n")
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        out.append(t)
            if sh.has_table:
                rows = []
                for row in sh.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    out.append("\n".join(rows))
        out.append("\n---")
    return "\n".join(out)


def write(name: str, content: str) -> None:
    target = DATA_DIR / name
    target.write_text(content, encoding="utf-8")
    print(f"  ✓ {name} ({len(content):,} chars)")


def copy(src_rel: str, dst_name: str) -> None:
    src = REPO_ROOT / src_rel
    if not src.exists():
        print(f"  ⚠ 來源不存在：{src_rel}")
        return
    shutil.copy2(src, DATA_DIR / dst_name)
    print(f"  ✓ {dst_name} (copied from {src_rel})")


def main() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    # 清掉舊輸出 — 避免重新命名/刪除 .pptx 後留下孤兒 .md 害 search_course 出現重複命中
    stale = list(DATA_DIR.glob("*.md"))
    for old in stale:
        old.unlink()
    print(f"→ writing to {DATA_DIR.relative_to(REPO_ROOT)} (cleared {len(stale)} stale .md)")

    # ── 抽 pptx ──
    pptx_to_md = {
        "01-why-mcp.pptx":              "01-why-mcp.md",
        "02-how-mcp-works.pptx":        "02-how-mcp-works.md",
        "03-agentic-tool-loop.pptx":    "03-agentic-tool-loop.md",
        "04-hands-on-lab.pptx":         "04-hands-on-lab-slides.md",
        "haiku-alignment-report.pptx":  "haiku-alignment-slides.md",
        "sonnet-running-example.pptx":  "sonnet-running-example-slides.md",
    }
    for src, dst in pptx_to_md.items():
        src_path = REPO_ROOT / src
        if src_path.exists():
            write(dst, extract_pptx(src_path))
        else:
            print(f"  ⚠ {src} 不存在，跳過")

    # ── 直接 cp 已是 md 的 ──
    md_copies = [
        ("04-hands-on-lab.md",                                       "04-hands-on-lab.md"),
        ("05-practical-considerations.md",                           "05-practical-considerations.md"),
        ("haiku-alignment-report.md",                                "haiku-alignment-report.md"),
        ("mini-project/docs/labs/L1-customize-your-data.md",         "lab-L1.md"),
        ("mini-project/docs/labs/L2-add-a-search-tool.md",           "lab-L2.md"),
        ("mini-project/docs/labs/L3-call-external-api.md",           "lab-L3.md"),
        ("mini-project/docs/benchmarks/2026-04-24-claude-vs-gemma4.md", "benchmark-claude-vs-gemma4.md"),
        ("mini-project/README.md",                                   "mini-project-readme.md"),
        ("course-notes-draft.md",                                    "course-notes-draft.md"),
    ]
    for src, dst in md_copies:
        copy(src, dst)

    print(f"\n✓ done. {len(list(DATA_DIR.glob('*.md')))} files in data/")


if __name__ == "__main__":
    main()
