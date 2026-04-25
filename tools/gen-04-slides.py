#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate 04-hands-on-lab.pptx from scratch.

Run:
    uv run --with python-pptx python3 tools/gen-04-slides.py

Output:
    ./04-hands-on-lab.pptx  (12 slides, 16:9)

Style conforms to the repo's visual spec (Ocean Gradient palette).
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── Palette (Ocean Gradient + accents) ──────────────────────────────
NAVY    = RGBColor(0x1E, 0x27, 0x61)
DEEP    = RGBColor(0x06, 0x5A, 0x82)
TEAL    = RGBColor(0x1C, 0x72, 0x93)
ORANGE  = RGBColor(0xE8, 0x79, 0x3A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CREAM   = RGBColor(0xFA, 0xF7, 0xF2)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
MUTED   = RGBColor(0x6B, 0x6B, 0x6B)
SOFT    = RGBColor(0xE8, 0xED, 0xF3)
CODE_BG = RGBColor(0x10, 0x1A, 0x2E)
CODE_FG = RGBColor(0xDD, 0xE4, 0xEE)

FONT_TITLE = 'Arial Black'
FONT_BODY  = 'Calibri'
FONT_CODE  = 'Consolas'

# ── Presentation setup ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def blank_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    # send to back
    spTree = s.shapes._spTree
    spTree.remove(rect._element); spTree.insert(2, rect._element)
    return s


def add_rect(slide, x, y, w, h, color, line_color=None):
    """Add filled rectangle at (x,y)  inches, (w,h) inches."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, text, *,
             font=FONT_BODY, size=18, color=DARK,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def add_multi(slide, x, y, w, h, paragraphs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a textbox with multiple paragraphs.

    Each element in `paragraphs` is a dict:
      { text, font, size, color, bold, italic, space_after }
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get('align', align)
        if 'space_after' in spec:
            p.space_after = Pt(spec['space_after'])
        r = p.add_run()
        r.text = spec['text']
        r.font.name = spec.get('font', FONT_BODY)
        r.font.size = Pt(spec.get('size', 18))
        r.font.color.rgb = spec.get('color', DARK)
        r.font.bold = spec.get('bold', False)
        r.font.italic = spec.get('italic', False)
    return tb


def page_chrome(slide, page_num, total, title, subtitle=None):
    """Shared chrome: Navy top bar, title, page number, orange accent."""
    # Navy left accent strip
    add_rect(slide, 0, 0, 0.35, 7.5, NAVY)
    # Title
    add_text(slide, 0.75, 0.35, 11.5, 0.8, title,
             font=FONT_TITLE, size=30, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, 0.75, 1.05, 11.5, 0.4, subtitle,
                 font=FONT_BODY, size=16, color=TEAL, italic=True)
    # Page number (right)
    add_text(slide, 11.8, 6.95, 1.4, 0.4, f'{page_num} / {total}',
             font=FONT_BODY, size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    # Orange thin line under title
    add_rect(slide, 0.75, 1.55, 3, 0.04, ORANGE)


TOTAL = 15


# ── Slide 1: Cover ──────────────────────────────────────────────────
def slide_1_cover():
    s = blank_slide(NAVY)
    # Orange accent square (design flourish)
    add_rect(s, 1.0, 2.3, 0.5, 0.5, ORANGE)
    # Small kicker
    add_text(s, 1.6, 2.35, 5, 0.5, 'Segment 4',
             font=FONT_BODY, size=20, color=ORANGE, bold=True)
    # Main title (Chinese - Calibri for CJK support)
    add_text(s, 1.0, 3.0, 11, 1.5, '動手做',
             font=FONT_TITLE, size=80, color=WHITE, bold=True)
    # Subtitle
    add_text(s, 1.05, 4.3, 11, 0.6, 'Hands-on Lab · 50 minutes',
             font=FONT_BODY, size=26, color=SOFT)
    # Separator
    add_rect(s, 1.0, 5.1, 2.5, 0.04, ORANGE)
    # Tagline
    add_text(s, 1.0, 5.25, 11, 0.5, 'mini-project 實作主場 — 現場跑起你自己領域的 AI agent',
             font=FONT_BODY, size=18, color=SOFT)
    # Footer
    add_text(s, 1.0, 6.85, 6, 0.4, 'NCHU MCP Workshop 2026',
             font=FONT_BODY, size=12, color=TEAL, italic=True)
    add_text(s, 7.0, 6.85, 5.8, 0.4, 'github.com/UDICatNCHU/nchu-mcp-workshop-2026',
             font=FONT_CODE, size=11, color=TEAL, align=PP_ALIGN.RIGHT, italic=True)


# ── Slide 2: 本節產出物 ──────────────────────────────────────────────
def slide_2_outcomes():
    s = blank_slide(WHITE)
    page_chrome(s, 2, TOTAL, '本節產出物',
                subtitle='50 分鐘後，你將親手擁有')

    rows = [
        ('✓', '在你自己電腦上跑起來的 MCP agent',
              'Node.js + Python FastMCP + 極簡 HTML chat，三層完整可運作'),
        ('✓', '一支「屬於你領域」的客製工具',
              '換 JSON 就好，0 行 Python；實驗室介紹 / 課程大綱 / 研究成果清單都行'),
        ('✓', '對 agent 資料流的「實際脈動感」',
              '工具選擇 → 參數綁定 → LLM 摘要這條路，你會在 terminal 看它一次跑通'),
        ('✓', 'L2 / L3 的明確挑戰路徑',
              '帶回去繼續深入：加有參數的搜尋工具、呼叫外部 API'),
    ]
    y = 2.0
    for mark, title, desc in rows:
        # Orange check
        add_text(s, 0.9, y, 0.7, 0.6, mark,
                 font=FONT_TITLE, size=32, color=ORANGE, bold=True)
        # Title + description
        add_multi(s, 1.5, y, 11.3, 1.3, [
            dict(text=title, font=FONT_BODY, size=21, color=DARK, bold=True, space_after=4),
            dict(text=desc,  font=FONT_BODY, size=15, color=MUTED),
        ])
        y += 1.15


# ── Slide 3: 時間配置 ──────────────────────────────────────────────
def slide_3_schedule():
    s = blank_slide(WHITE)
    page_chrome(s, 3, TOTAL, '50 分鐘時間配置',
                subtitle='講師 10 min 引導 + 40 min 巡場陪跑')

    rows = [
        ('0–10',  '講師 demo + 學員同步 ./setup.sh',       '環境綠燈 5/5 ✅',  DEEP),
        ('10–20', 'L1 Step 1–2：觀察現況 + 換自己 JSON',    'data/your.json',   TEAL),
        ('20–35', 'L1 Step 3–4：改 docstring + 重啟驗證',  '問自己資料會答',   TEAL),
        ('35–45', '交叉展示：3–4 位老師 demo 自己領域',    '見識不同落地方式', DEEP),
        ('45–50', 'Q&A + 為 L2/L3 與 Segment 5 鋪陳',      '清楚下一步',       ORANGE),
    ]

    # Table header row
    header_y = 2.0
    add_rect(s, 0.75, header_y, 12.0, 0.55, NAVY)
    add_text(s, 0.95, header_y + 0.08, 1.7, 0.4, '時段 (min)',
             font=FONT_BODY, size=14, color=WHITE, bold=True)
    add_text(s, 2.85, header_y + 0.08, 6.5, 0.4, '做什麼',
             font=FONT_BODY, size=14, color=WHITE, bold=True)
    add_text(s, 9.5,  header_y + 0.08, 3.2, 0.4, '產出',
             font=FONT_BODY, size=14, color=WHITE, bold=True)

    y = 2.55
    for rng, what, out, color in rows:
        add_rect(s, 0.75, y, 12.0, 0.7, CREAM)
        add_rect(s, 0.75, y, 0.12, 0.7, color)  # left accent
        add_text(s, 0.95, y + 0.15, 1.7, 0.4, rng,
                 font=FONT_CODE, size=15, color=color, bold=True)
        add_text(s, 2.85, y + 0.15, 6.5, 0.4, what,
                 font=FONT_BODY, size=15, color=DARK)
        add_text(s, 9.5, y + 0.15, 3.2, 0.4, out,
                 font=FONT_BODY, size=14, color=MUTED, italic=True)
        y += 0.8

    # Footnote
    add_text(s, 0.75, 6.7, 12, 0.4,
             '※ 40 分鐘巡場預計 80% 在幫卡住的老師；開場預演越短，越多時間陪跑。',
             font=FONT_BODY, size=12, color=MUTED, italic=True)


# ── Slide 4: 架構回顧 ──────────────────────────────────────────────
def slide_4_architecture():
    s = blank_slide(WHITE)
    page_chrome(s, 4, TOTAL, '架構快速回顧',
                subtitle='不是單向 pipeline — Node ⟷ LLM 的多輪迭代才是 agent')

    boxes = [
        ('Browser', 'web/index.html  ·  fetch POST /chat', DEEP),
        ('Node server', 'Express  ·  LLMClient  ·  MCPClient', TEAL),
        ('Python FastMCP', '@mcp.tool()  ·  stdio JSON-RPC', DEEP),
        ('data/', '你的 JSON  ·  外部 API', NAVY),
    ]
    x0, y0 = 1.0, 2.0
    w, h = 7.6, 0.95
    for i, (title, desc, color) in enumerate(boxes):
        y = y0 + i * (h + 0.2)
        add_rect(s, x0, y, w, h, CREAM)
        add_rect(s, x0, y, 0.15, h, color)
        add_text(s, x0 + 0.35, y + 0.1, 3.5, 0.4, title,
                 font=FONT_TITLE, size=18, color=color, bold=True)
        add_text(s, x0 + 0.35, y + 0.5, w - 0.5, 0.4, desc,
                 font=FONT_CODE, size=13, color=DARK)
        # Bidirectional arrow between boxes
        if i < len(boxes) - 1:
            add_text(s, x0 + w/2 - 0.3, y + h - 0.05, 0.6, 0.3, '⇅',
                     font=FONT_TITLE, size=22, color=ORANGE,
                     bold=True, align=PP_ALIGN.CENTER)

    # LLM API side bubble at Node row (i=1)
    llm_y = y0 + 1 * (h + 0.2)
    bub_x = x0 + w + 0.7
    bub_w = 3.6
    add_rect(s, bub_x, llm_y, bub_w, h, NAVY)
    add_text(s, bub_x + 0.15, llm_y + 0.1, bub_w - 0.3, 0.4, '☁  LLM API',
             font=FONT_TITLE, size=17, color=WHITE, bold=True)
    add_text(s, bub_x + 0.15, llm_y + 0.5, bub_w - 0.3, 0.4,
             'Claude · Gemma · GPT · ...',
             font=FONT_CODE, size=13, color=SOFT)
    # Node ⟷ LLM arrow
    add_text(s, x0 + w + 0.05, llm_y + 0.25, 0.65, 0.45, '⟷',
             font=FONT_TITLE, size=24, color=ORANGE,
             bold=True, align=PP_ALIGN.CENTER)
    # iteration label below bubble
    add_text(s, bub_x, llm_y + h + 0.05, bub_w, 0.3, 'agent loop · 多輪迭代',
             font=FONT_BODY, size=12, color=ORANGE,
             italic=True, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, 0.75, 6.85, 12, 0.4,
             '關鍵：⇅ 雙向 — Node ⟷ LLM 之間迭代「問 → tool_use → 跑 → 答 → 再問」直到 end_turn',
             font=FONT_BODY, size=13, color=MUTED, italic=True)


# ── Slide 5: 20 行的靈魂（agent loop 真實程式碼）──────────────────
def slide_4b_agent_loop():
    s = blank_slide(WHITE)
    page_chrome(s, 5, TOTAL, 'agent loop 真實長這樣',
                subtitle='backend-node/llm-client.js — 整個系統的 20 行核心')

    # Code block (left ~ 7.7 inch)
    add_rect(s, 0.75, 1.95, 7.7, 4.95, CODE_BG)
    code = (
        'async chat(messages) {\n'
        '  const history = [...messages];\n'
        '  for (let i = 0; i < maxIterations; i++) {           // ← 護欄\n'
        '    const resp = await anthropic.messages.create({\n'
        '      model, tools: mcp.getAnthropicTools(),          // ① 餵 tools\n'
        '      messages: history,\n'
        '    });\n'
        '    history.push({ role: "assistant", content: resp.content });\n'
        '\n'
        '    if (resp.stop_reason !== "tool_use") {            // ② 結束\n'
        '      return { reply: extractText(resp), messages: history };\n'
        '    }\n'
        '\n'
        '    const toolUses = resp.content                     // ③ 跑工具\n'
        '      .filter(b => b.type === "tool_use");\n'
        '    const toolResults = await Promise.all(\n'
        '      toolUses.map(t => mcp.callTool(t.name, t.input))\n'
        '    );\n'
        '    history.push({ role: "user", content: toolResults });\n'
        '  }\n'
        '}'
    )
    add_text(s, 0.95, 2.05, 7.4, 4.8, code,
             font=FONT_CODE, size=12, color=CODE_FG)

    # Right side: 3 annotation cards
    cards = [
        ('①', '餵', 'messages + tools 全部重餵 — \nAPI 是 stateless，每輪都要全 history'),
        ('②', '結束', 'stop_reason ≠ tool_use\n→ 取 text，回給前端'),
        ('③', '工具', 'tool_use → 平行 callTool\n→ 結果 push 回 history\n→ 繼續迴圈'),
    ]
    x0, y0 = 8.65, 1.95
    for i, (num, title, desc) in enumerate(cards):
        y = y0 + i * 1.4
        add_rect(s, x0, y, 4.1, 1.25, CREAM)
        add_rect(s, x0, y, 0.5, 1.25, ORANGE)
        add_text(s, x0 + 0.05, y + 0.32, 0.5, 0.6, num,
                 font=FONT_TITLE, size=22, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, x0 + 0.65, y + 0.08, 3.4, 0.4, title,
                 font=FONT_BODY, size=15, color=NAVY, bold=True)
        add_text(s, x0 + 0.65, y + 0.42, 3.4, 0.85, desc,
                 font=FONT_BODY, size=11, color=DARK)

    add_text(s, 0.75, 6.95, 12, 0.4,
             '🛑 maxIterations=10 是護欄；超過 → 拋錯。整個 agent 沒有 magic — 就是個 for-loop + if-else。',
             font=FONT_BODY, size=12, color=ORANGE, bold=True, italic=True)


# ── Slide 6: 行前準備 ──────────────────────────────────────────────
def slide_5_pre_workshop():
    s = blank_slide(WHITE)
    page_chrome(s, 6, TOTAL, '行前準備（請於上課前完成）',
                subtitle='這幾步無關概念，但卡住會吃掉現場時間')

    rows = [
        ('1', 'git clone 教材並進到 mini-project',
              '$ git clone https://github.com/UDICatNCHU/nchu-mcp-workshop-2026\n'
              '$ cd nchu-mcp-workshop-2026/mini-project'),
        ('2', '取得 LLM 存取（二選一）',
              '雲端：Anthropic Console 申請 API key（新帳號送 $5 試用）\n'
              '本地：確認你會走 NCHU vLLM 路線（細節見下一頁）'),
        ('3', '建 .env 並填入金鑰',
              '$ cp .env.example .env\n'
              '$ vim .env   # 填 ANTHROPIC_API_KEY 或 OPENAI_BASE_URL'),
        ('4', '跑環境預檢',
              '$ ./setup.sh\n'
              '    → 看到 5/5 ✅ 代表現場可以直接 npm start'),
    ]
    y = 1.95
    for num, title, code in rows:
        add_rect(s, 0.75, y, 12.0, 1.05, CREAM)
        add_rect(s, 0.75, y, 0.55, 1.05, DEEP)
        add_text(s, 0.78, y + 0.28, 0.55, 0.5, num,
                 font=FONT_TITLE, size=26, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, 1.5, y + 0.08, 11.2, 0.4, title,
                 font=FONT_BODY, size=16, color=DARK, bold=True)
        add_text(s, 1.5, y + 0.45, 11.2, 0.6, code,
                 font=FONT_CODE, size=12, color=DEEP)
        y += 1.18

    add_text(s, 0.75, 6.7, 12, 0.4,
             '✅ 行前 4 步完成 → 現場只剩「npm start + 開瀏覽器」（下一頁）',
             font=FONT_BODY, size=13, color=ORANGE, bold=True, italic=True)


# ── Slide 7: 現場 Quick Start ────────────────────────────────────────
def slide_5_quickstart():
    s = blank_slide(WHITE)
    page_chrome(s, 7, TOTAL, '現場啟動（5 分鐘）',
                subtitle='三條指令 + 對照「該長什麼樣」')

    # Section 1: commands you type
    add_text(s, 0.75, 1.9, 12, 0.35, '▶  你輸入',
             font=FONT_BODY, size=14, color=NAVY, bold=True)
    add_rect(s, 0.75, 2.25, 12.0, 1.55, CODE_BG)
    cmds = (
        '$ ./setup.sh                            # 1. sanity check 已裝好的環境\n'
        '$ cd backend-node && npm start          # 2. 啟動 backend + spawn MCP\n'
        '$ open http://localhost:3000            # 3. 瀏覽器'
    )
    add_text(s, 1.0, 2.35, 11.5, 1.4, cmds,
             font=FONT_CODE, size=14, color=CODE_FG)

    # Section 2: what you should see
    add_text(s, 0.75, 3.95, 12, 0.35, '▶  你應該看到（npm start 印出）',
             font=FONT_BODY, size=14, color=NAVY, bold=True)
    add_rect(s, 0.75, 4.3, 12.0, 2.05, CODE_BG)
    expected = (
        '> mini-assistant@1.0.0 start\n'
        '> node server.js\n'
        '\n'
        '✓ hello_tool → get_english_center_info\n'
        '✓ teachers_tool → search_teachers, get_teacher_detail\n'
        '✓ weather_tool → get_weather\n'
        '→ Mini AI Assistant: http://localhost:3000'
    )
    add_text(s, 1.0, 4.4, 11.5, 1.9, expected,
             font=FONT_CODE, size=12, color=CODE_FG)

    add_text(s, 0.75, 6.5, 12, 0.4,
             '🎉 三行 ✓ + URL → MCP server 都連上了；瀏覽器問「英文中心幾點開門？」會在 terminal 看 [tool_use]',
             font=FONT_BODY, size=12, color=ORANGE, bold=True, italic=True)
    add_text(s, 0.75, 6.85, 12, 0.4,
             '⚠ 沒看到 ✓ 三行：→ Slide 13 卡點速查',
             font=FONT_BODY, size=11, color=MUTED, italic=True)


# ── Slide 6: L1 Overview ──────────────────────────────────────────
def slide_6_l1_overview():
    s = blank_slide(WHITE)
    page_chrome(s, 9, TOTAL, 'Lab 1 — 換 JSON 做你領域的助理',
                subtitle='0 行 Python · 40 分鐘 · 現場必做')

    # Big 4-step flow
    steps = [
        ('1', '觀察', '看 [tool_use] log\nAI 摘要 JSON',          '5 min',  DEEP),
        ('2', '換資料', '編輯 data/*.json\n放你的領域內容',       '15 min', TEAL),
        ('3', '改說明書', 'docstring 的使用情境\n決定工具呼叫率', '10 min', ORANGE),
        ('4', '驗證', '重啟 npm start\n問自己資料問題',           '5 min',  DEEP),
    ]
    x0, y = 0.85, 2.3
    box_w, box_h = 2.9, 3.8
    gap = 0.15
    for i, (num, title, desc, dur, color) in enumerate(steps):
        x = x0 + i * (box_w + gap)
        add_rect(s, x, y, box_w, box_h, CREAM)
        add_rect(s, x, y, box_w, 0.75, color)
        # Number (white on color)
        add_text(s, x + 0.3, y + 0.08, 1.0, 0.6, num,
                 font=FONT_TITLE, size=36, color=WHITE, bold=True)
        # Step title
        add_text(s, x + 1.1, y + 0.15, box_w - 1.2, 0.5, title,
                 font=FONT_TITLE, size=22, color=WHITE, bold=True)
        # Description
        add_text(s, x + 0.25, y + 1.0, box_w - 0.5, 2.0, desc,
                 font=FONT_BODY, size=15, color=DARK)
        # Duration badge
        add_rect(s, x + 0.25, y + 3.15, 1.5, 0.45, color)
        add_text(s, x + 0.25, y + 3.2, 1.5, 0.4, dur,
                 font=FONT_CODE, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, 0.85, 6.6, 12, 0.5,
             '產出：你的 AI 助理會回答「你自己放進去的資料」— 不是英文中心 demo 了',
             font=FONT_BODY, size=15, color=ORANGE, bold=True, italic=True)


# ── Slide 7: Step 1-2: 觀察 + 換 JSON ──────────────────────────────
def slide_7_step12():
    s = blank_slide(WHITE)
    page_chrome(s, 10, TOTAL, 'Step 1–2：觀察 → 換 JSON',
                subtitle='還沒動到 Python 一個字')

    # Left: Step 1 观察
    add_rect(s, 0.75, 2.0, 6.0, 4.6, CREAM)
    add_rect(s, 0.75, 2.0, 6.0, 0.55, DEEP)
    add_text(s, 1.0, 2.07, 5.8, 0.45, 'Step 1　觀察（5 min）',
             font=FONT_TITLE, size=18, color=WHITE, bold=True)

    add_multi(s, 0.95, 2.75, 5.7, 3.8, [
        dict(text='在瀏覽器問：', size=15, color=DARK, space_after=2),
        dict(text='  「英文中心幾點開門？」', font=FONT_CODE, size=14, color=DEEP, space_after=8),
        dict(text='在 terminal 觀察：', size=15, color=DARK, space_after=2),
        dict(text='  [tool_use] get_english_center_info', font=FONT_CODE, size=13, color=ORANGE, space_after=10),
        dict(text='💡 關鍵觀察', size=15, color=DARK, bold=True, space_after=2),
        dict(text='Claude 不會 dump 整份 JSON。它挑出「開放時間」那一段。',
             size=13, color=MUTED, space_after=4),
        dict(text='這是 LLM 對 tool result 做的「摘要」— agent 的靈魂。',
             size=13, color=MUTED),
    ])

    # Right: Step 2 換 JSON
    add_rect(s, 7.0, 2.0, 5.75, 4.6, CREAM)
    add_rect(s, 7.0, 2.0, 5.75, 0.55, TEAL)
    add_text(s, 7.2, 2.07, 5.5, 0.45, 'Step 2　換 JSON（15 min）',
             font=FONT_TITLE, size=18, color=WHITE, bold=True)

    add_multi(s, 7.15, 2.75, 5.55, 3.8, [
        dict(text='編輯：', size=14, color=DARK, space_after=2),
        dict(text='  mcp-server-py/data/english_center.json',
             font=FONT_CODE, size=12, color=TEAL, space_after=8),
        dict(text='換成你自己的領域資料：', size=14, color=DARK, space_after=2),
        dict(text='• 你的研究室介紹', size=13, color=DARK, space_after=1),
        dict(text='• 你教的一門課（大綱/作業/評分）', size=13, color=DARK, space_after=1),
        dict(text='• 你系所的 FAQ', size=13, color=DARK, space_after=1),
        dict(text='• 你的研究成果清單', size=13, color=DARK, space_after=8),
        dict(text='⚠ JSON 語法錯會卡住：', size=13, color=ORANGE, bold=True, space_after=2),
        dict(text='  python3 -m json.tool data/xxx.json  驗一下',
             font=FONT_CODE, size=12, color=MUTED),
    ])


# ── Slide 8: Step 3 Docstring ──────────────────────────────────────
def slide_8_step3():
    s = blank_slide(WHITE)
    page_chrome(s, 11, TOTAL, 'Step 3：改 docstring（10 min）',
                subtitle='直接拿 repo 兩支真實工具對照 — 同樣語法、教 LLM 不同事')

    # Left: minimal (hello_tool.py)
    add_rect(s, 0.75, 1.95, 6.0, 4.5, CREAM)
    add_rect(s, 0.75, 1.95, 6.0, 0.5, DEEP)
    add_text(s, 1.0, 2.0, 5.8, 0.4, '📝  minimal — hello_tool.py',
             font=FONT_TITLE, size=15, color=WHITE, bold=True)
    minimal = (
        '@mcp.tool()\n'
        'def get_english_center_info() -> str:\n'
        '    """取得中興大學英語自學暨\n'
        '    檢定中心的完整資訊。\n'
        '\n'
        '    回傳 JSON 字串，包含：\n'
        '    名稱、開放時間、地點、設備…\n'
        '    使用情境：使用者詢問英語\n'
        '    自學中心相關問題時呼叫。\n'
        '    """'
    )
    add_text(s, 0.9, 2.55, 5.8, 3.85, minimal,
             font=FONT_CODE, size=12, color=DARK)

    # Right: production-grade (teachers_tool.py - get_teacher_detail)
    add_rect(s, 7.0, 1.95, 5.75, 4.5, CREAM)
    add_rect(s, 7.0, 1.95, 5.75, 0.5, ORANGE)
    add_text(s, 7.2, 2.0, 5.5, 0.4, '📖  rich — teachers_tool.py',
             font=FONT_TITLE, size=15, color=WHITE, bold=True)
    rich = (
        '@mcp.tool()\n'
        'def get_teacher_detail(name: str) -> str:\n'
        '    """取得指定教授的完整資訊\n'
        '    （email、辦公室、研究領域）。\n'
        '\n'
        '    使用情境：使用者問某位教授\n'
        '    的聯絡方式或完整資料時呼叫。\n'
        '    通常在 search_teachers 找到\n'
        '    候選名單後再呼叫。 ←跨工具線索\n'
        '\n'
        '    Args:\n'
        '        name: 教授姓名（完整中文）。\n'
        '    """'
    )
    add_text(s, 7.15, 2.55, 5.55, 3.85, rich,
             font=FONT_CODE, size=11, color=DARK)

    # Footer: 3 key elements
    add_rect(s, 0.75, 6.55, 12.0, 0.55, NAVY)
    add_text(s, 0.95, 6.6, 11.7, 0.5,
             '影響 LLM 行為的 3 個元素：'
             ' ① 使用情境 → 該不該叫'
             '    ② Args 描述 → 怎麼填參數'
             '    ③ 跨工具線索 → 多輪呼叫順序',
             font=FONT_BODY, size=12, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)


# ── Slide 9: Step 4 驗證 ──────────────────────────────────────────
def slide_9_step4():
    s = blank_slide(WHITE)
    page_chrome(s, 12, TOTAL, 'Step 4：重啟 + 驗證（5 min）',
                subtitle='三件事證明你的工具「真的在 work」')

    # Command block
    add_rect(s, 0.75, 2.0, 12.0, 1.3, CODE_BG)
    add_text(s, 1.0, 2.1, 11.5, 1.1,
             '$ Ctrl+C                            # 停掉前一次 server\n'
             '$ cd backend-node && npm start\n'
             '  → ✓ hello_tool → get_lab_info    # 新名字出現代表工具註冊成功',
             font=FONT_CODE, size=15, color=CODE_FG)

    # Three verify cards
    cards = [
        ('①', '問「不相關」的問題',     '「今天天氣？」',  'AI 應該不呼叫你的工具，直接聊天',  DEEP),
        ('②', '問「相關」的問題',       '「研究室招生條件？」',  'AI 會引用 JSON 裡的實際內容',  TEAL),
        ('③', '問「資料沒有」的內容',   '「有實習機會嗎？」（JSON 沒這欄）',
              'AI 應該明說不知道，不會幻想',  ORANGE),
    ]
    y = 3.65
    for num, title, q, expect, color in cards:
        add_rect(s, 0.75, y, 12.0, 0.75, CREAM)
        add_rect(s, 0.75, y, 0.12, 0.75, color)
        add_text(s, 1.0, y + 0.12, 0.5, 0.5, num,
                 font=FONT_TITLE, size=22, color=color, bold=True)
        add_text(s, 1.6, y + 0.06, 4.2, 0.4, title,
                 font=FONT_BODY, size=15, color=DARK, bold=True)
        add_text(s, 1.6, y + 0.4, 4.2, 0.4, q,
                 font=FONT_CODE, size=13, color=MUTED)
        add_text(s, 5.9, y + 0.18, 6.9, 0.5, expect,
                 font=FONT_BODY, size=14, color=DARK)
        y += 0.9

    add_text(s, 0.75, 6.65, 12, 0.4,
             '第 ③ 項過關最關鍵 — 代表 agent 在 tool-grounded，不是在胡謅。',
             font=FONT_BODY, size=13, color=MUTED, italic=True)


# ── Slide 8: 選你的 LLM 路線 ──────────────────────────────────────
def slide_10_local_endpoint():
    s = blank_slide(WHITE)
    page_chrome(s, 8, TOTAL, '選你的 LLM 路線',
                subtitle='`.env` 一行切換 — 因為 MCP 工具兩家都認')

    # Top row: two .env columns (compressed)
    add_rect(s, 0.75, 1.95, 6.0, 1.95, CREAM)
    add_rect(s, 0.75, 1.95, 6.0, 0.5, DEEP)
    add_text(s, 1.0, 2.0, 5.8, 0.4, '☁  雲端 · Anthropic Claude',
             font=FONT_TITLE, size=15, color=WHITE, bold=True)
    add_rect(s, 0.95, 2.55, 5.6, 1.3, CODE_BG)
    add_text(s, 1.05, 2.62, 5.45, 1.2,
             '# .env\n'
             'LLM_PROVIDER=anthropic\n'
             'ANTHROPIC_API_KEY=sk-ant-...\n'
             'ANTHROPIC_MODEL=claude-haiku-4-5',
             font=FONT_CODE, size=12, color=CODE_FG)

    add_rect(s, 7.0, 1.95, 5.75, 1.95, CREAM)
    add_rect(s, 7.0, 1.95, 5.75, 0.5, TEAL)
    add_text(s, 7.2, 2.0, 5.5, 0.4, '🏠  本地 · NCHU vLLM (Gemma 4)',
             font=FONT_TITLE, size=15, color=WHITE, bold=True)
    add_rect(s, 7.15, 2.55, 5.45, 1.3, CODE_BG)
    add_text(s, 7.25, 2.62, 5.3, 1.2,
             '# .env\n'
             'LLM_PROVIDER=openai\n'
             'OPENAI_BASE_URL=http://<ws-host>:8000/v1\n'
             'OPENAI_API_KEY=dummy\n'
             'OPENAI_MODEL=gemma-4',
             font=FONT_CODE, size=12, color=CODE_FG)

    # MIDDLE: N+M proof — same MCP tool, two adapter formats (4+6 lines each)
    add_rect(s, 0.75, 4.0, 12.0, 1.85, CREAM)
    add_rect(s, 0.75, 4.0, 0.12, 1.85, ORANGE)
    add_text(s, 1.0, 4.05, 11.7, 0.35,
             '💎  N+M 的甜蜜：你的工具一支不變，adapter 差別只在 mcp-client.js 這幾行',
             font=FONT_BODY, size=13, color=NAVY, bold=True)

    # Anthropic format (left)
    add_rect(s, 1.0, 4.45, 5.5, 1.3, CODE_BG)
    add_text(s, 1.1, 4.5, 5.35, 1.2,
             '// getAnthropicTools()\n'
             '{\n'
             '  name,\n'
             '  description,\n'
             '  input_schema,   // ← MCP\'s\n'
             '}',
             font=FONT_CODE, size=11, color=CODE_FG)

    # OpenAI format (right)
    add_rect(s, 7.0, 4.45, 5.5, 1.3, CODE_BG)
    add_text(s, 7.1, 4.5, 5.35, 1.2,
             '// getOpenAITools()\n'
             '{ type: "function",\n'
             '  function: {\n'
             '    name, description,\n'
             '    parameters,   // ← 同一個\n'
             '} }',
             font=FONT_CODE, size=11, color=CODE_FG)

    # Bottom: benchmark strip
    add_rect(s, 0.75, 5.95, 12.0, 0.55, NAVY)
    add_text(s, 0.95, 5.98, 11.7, 0.5,
             '📊 2026-04-24 實測：Tool 正確率 Claude 100% · Gemma 4 100%　|　'
             '延遲 4.8s vs 3.5s　|　兩者 hallucination 都通過',
             font=FONT_CODE, size=11, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 0.75, 6.6, 12, 0.4,
             '→ 詳細 benchmark：mini-project/docs/benchmarks/2026-04-24-claude-vs-gemma4.md',
             font=FONT_CODE, size=11, color=MUTED, italic=True)


# ── Slide 11: 常見卡點 ────────────────────────────────────────────
def slide_11_pitfalls():
    s = blank_slide(WHITE)
    page_chrome(s, 13, TOTAL, '常見卡點速查',
                subtitle='卡住 30 秒就看這頁；更多在 Lab 手冊結尾')

    rows = [
        ('setup.sh 卡 Node 版本',           '本機 Node < 18',            'nvm install 22 && nvm use 22'),
        ('port 3000 已被占',                 '其他服務先佔了',            'PORT=3001 npm start'),
        ('瀏覽器問後沒反應',                 '.env 的 API key 還是 placeholder', '檢查 sk-ant-... 有沒有填'),
        ('JSON 讀取錯誤',                    '手改 JSON 有語法錯',        'python3 -m json.tool data/xxx.json'),
        ('問相關問題但 LLM 不叫工具',        'docstring 寫太短 / 太含糊', '補「使用情境：當使用者問 X 時」'),
        ('npm start 看不到工具列表',         'config.json 的 server key\n'
                                             '與 FastMCP("...") 名稱對不上', '兩邊改成同一個字串'),
        ('Could not find tool / 啟動失敗',  'mcp-server-py 沒裝依賴',    'cd mcp-server-py && uv sync'),
    ]

    # Header
    add_rect(s, 0.75, 1.85, 12.0, 0.45, NAVY)
    add_text(s, 0.95, 1.9, 3.9, 0.35, '症狀',
             font=FONT_BODY, size=13, color=WHITE, bold=True)
    add_text(s, 5.0,  1.9, 3.7, 0.35, '原因',
             font=FONT_BODY, size=13, color=WHITE, bold=True)
    add_text(s, 8.8,  1.9, 4.0, 0.35, '解法',
             font=FONT_BODY, size=13, color=WHITE, bold=True)

    y = 2.32
    for sym, cause, fix in rows:
        add_rect(s, 0.75, y, 12.0, 0.62, CREAM)
        add_text(s, 0.95, y + 0.1, 3.9, 0.5, sym,
                 font=FONT_BODY, size=12, color=DARK)
        add_text(s, 5.0,  y + 0.1, 3.7, 0.5, cause,
                 font=FONT_BODY, size=11, color=MUTED)
        add_text(s, 8.8,  y + 0.1, 4.0, 0.5, fix,
                 font=FONT_CODE, size=11, color=DEEP)
        y += 0.66

    add_text(s, 0.75, 6.95, 12, 0.4,
             '卡住超過 3 分鐘：舉手。講師 + 鄰座老師會過來幫你。',
             font=FONT_BODY, size=12, color=ORANGE, bold=True, italic=True)


# ── Slide 13: Show & Tell（交叉展示） ─────────────────────────────
def slide_13_show_and_tell():
    s = blank_slide(WHITE)
    page_chrome(s, 14, TOTAL, '交叉展示 — Show & Tell',
                subtitle='35–45 min · 看 3 位老師把同一支 mini-project 變成自己的 agent')

    # Left: 講師 prompt
    add_rect(s, 0.75, 2.0, 6.0, 4.6, CREAM)
    add_rect(s, 0.75, 2.0, 6.0, 0.55, ORANGE)
    add_text(s, 1.0, 2.07, 5.8, 0.45, '🎤  請 3 位老師上來（每人 3 分鐘）',
             font=FONT_TITLE, size=16, color=WHITE, bold=True)

    add_multi(s, 0.95, 2.8, 5.7, 3.7, [
        dict(text='① 一句話介紹你的領域 + 你放了什麼 JSON',
             font=FONT_BODY, size=15, color=DARK, bold=True, space_after=4),
        dict(text='   例：「我教植物分類學，JSON 是 50 種校園樹木」',
             font=FONT_BODY, size=12, color=MUTED, italic=True, space_after=12),
        dict(text='② 問你的 agent 一個你自己關心的問題',
             font=FONT_BODY, size=15, color=DARK, bold=True, space_after=4),
        dict(text='   例：「校門口那棵開白花的是什麼？」',
             font=FONT_BODY, size=12, color=MUTED, italic=True, space_after=12),
        dict(text='③ 講一個你踩到的雷（30 秒）',
             font=FONT_BODY, size=15, color=DARK, bold=True, space_after=4),
        dict(text='   docstring 寫太短？JSON 結構太深？編碼錯？',
             font=FONT_BODY, size=12, color=MUTED, italic=True),
    ])

    # Right: 為什麼做這件事
    add_rect(s, 7.0, 2.0, 5.75, 4.6, CREAM)
    add_rect(s, 7.0, 2.0, 5.75, 0.55, DEEP)
    add_text(s, 7.2, 2.07, 5.5, 0.45, '💡  為什麼花 10 分鐘做這件事',
             font=FONT_TITLE, size=16, color=WHITE, bold=True)

    add_multi(s, 7.15, 2.8, 5.55, 3.7, [
        dict(text='看見落地的多樣性',
             font=FONT_BODY, size=15, color=DARK, bold=True, space_after=4),
        dict(text='醫學人文 / 物理 / 語言學 / 行政 FAQ —\n'
                  '同一支 code 在不同領域的形狀。',
             font=FONT_BODY, size=12, color=MUTED, space_after=12),
        dict(text='偷學別人 docstring 寫法',
             font=FONT_BODY, size=15, color=DARK, bold=True, space_after=4),
        dict(text='10 個老師有 10 種「使用情境」描述法 —\n'
                  '看誰的 AI 最聽話。',
             font=FONT_BODY, size=12, color=MUTED, space_after=12),
        dict(text='提早為 L2 / L3 預警',
             font=FONT_BODY, size=15, color=DARK, bold=True, space_after=4),
        dict(text='別人卡的雷，可能就是你下午自修會踩到的。',
             font=FONT_BODY, size=12, color=MUTED),
    ])

    add_text(s, 0.75, 6.85, 12, 0.4,
             '※ 沒上台的老師不是觀眾 — 把對你「最有啟發」的那個 demo 記下來，課後試試',
             font=FONT_BODY, size=12, color=ORANGE, italic=True, bold=True)


# ── Slide 15: 結束 & 鋪陳 Segment 5 ────────────────────────────────
def slide_12_end():
    s = blank_slide(NAVY)

    # Big congrats
    add_text(s, 0.75, 1.2, 12, 0.8, '恭喜 — 你剛做出了',
             font=FONT_BODY, size=26, color=SOFT)
    add_text(s, 0.75, 1.9, 12, 1.2, '屬於你自己領域的 AI agent',
             font=FONT_TITLE, size=48, color=WHITE, bold=True)
    add_rect(s, 0.75, 3.1, 2.5, 0.05, ORANGE)

    # Two columns
    # Left: 回家路徑
    add_text(s, 0.75, 3.45, 6, 0.4, '回家路徑（課後自修）',
             font=FONT_TITLE, size=18, color=ORANGE, bold=True)
    add_multi(s, 0.75, 3.9, 6, 2.5, [
        dict(text='🧪  L2 — 加一支有參數的搜尋工具',
             font=FONT_BODY, size=16, color=WHITE, bold=True, space_after=4),
        dict(text='     docs/labs/L2-add-a-search-tool.md  (40 分)',
             font=FONT_CODE, size=12, color=SOFT, space_after=10),
        dict(text='🧪  L3 — 呼叫外部 API（async + XML）',
             font=FONT_BODY, size=16, color=WHITE, bold=True, space_after=4),
        dict(text='     docs/labs/L3-call-external-api.md  (60 分)',
             font=FONT_CODE, size=12, color=SOFT),
    ])

    # Right: 鋪陳 Segment 5
    add_text(s, 7.0, 3.45, 5.8, 0.4, '下一段 Segment 5',
             font=FONT_TITLE, size=18, color=ORANGE, bold=True)
    add_multi(s, 7.0, 3.9, 5.8, 2.5, [
        dict(text='實務考量（10 min 收尾）',
             font=FONT_BODY, size=16, color=WHITE, bold=True, space_after=6),
        dict(text='從「你剛做的 3 個工具」往外擴：',
             font=FONT_BODY, size=14, color=SOFT, space_after=3),
        dict(text='• 239 工具怎麼辦？（Scale）',
             font=FONT_BODY, size=13, color=SOFT, space_after=2),
        dict(text='• 工具描述怎麼寫？（Quality）',
             font=FONT_BODY, size=13, color=SOFT, space_after=2),
        dict(text='• 該用哪個模型？（Model）',
             font=FONT_BODY, size=13, color=SOFT, space_after=2),
        dict(text='• 帳單會有多貴？（Cost）',
             font=FONT_BODY, size=13, color=SOFT),
    ])

    # Footer
    add_text(s, 0.75, 6.7, 12, 0.4,
             'github.com/UDICatNCHU/nchu-mcp-workshop-2026   ·   Issues 歡迎！',
             font=FONT_CODE, size=12, color=TEAL, align=PP_ALIGN.CENTER, italic=True)


# ── Build ───────────────────────────────────────────────────────────
BUILDERS = [
    slide_1_cover,           # 1  Cover
    slide_2_outcomes,        # 2  本節產出物
    slide_3_schedule,        # 3  時間配置
    slide_4_architecture,    # 4  架構回顧
    slide_4b_agent_loop,     # 5  20 行靈魂（agent loop 真實程式碼）
    slide_5_pre_workshop,    # 6  行前準備
    slide_5_quickstart,      # 7  現場 Quick Start
    slide_10_local_endpoint, # 8  選你的 LLM 路線
    slide_6_l1_overview,     # 9  L1 Overview
    slide_7_step12,          # 10 Step 1-2
    slide_8_step3,           # 11 Step 3 docstring
    slide_9_step4,           # 12 Step 4 驗證
    slide_11_pitfalls,       # 13 卡點速查
    slide_13_show_and_tell,  # 14 Show & Tell
    slide_12_end,            # 15 結束
]
for b in BUILDERS:
    b()

out = Path(__file__).resolve().parent.parent / '04-hands-on-lab.pptx'
prs.save(out)
print(f'✓ wrote {out} ({len(prs.slides)} slides)')
